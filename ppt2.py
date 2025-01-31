import streamlit as st
from crewai import Agent, Task, Crew, LLM
from crewai_tools import SerperDevTool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
 
 
def create_presentation(slides_text, output_file="output_presentation.pptx"):
    prs = Presentation()
    slides_content = slides_text.split("Slide")
   
    for i, point in enumerate(slides_content):
        if point.strip():
            slide_layout = prs.slide_layouts[1]  # Title + Content Layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
 
            title.text = f"Slide {i+1}"
            content.text = point
 
            # Style Text
            for paragraph in content.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT  # Align text
                paragraph.font.bold = True  # Make text bold
                paragraph.font.size = Pt(24)  # Set font size
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color
 
    prs.save(output_file)
    return output_file
 
def main():
    st.title("Presentation Generator")
    serpapi_key = st.text_input("Enter your SerpAPI key:")
    openai_api_key = st.text_input("Enter your OpenAI API key:")
    text = st.text_input("Input the prompt for which you want to create a presentation:")
   
    if st.button("Generate Presentation"):
        if text:
           
            # Initialize the regular model
            llm = LLM(
                model="gpt-3.5-turbo-16k", temperature=0.3, api_key=openai_api_key
            )
            search_tool = SerperDevTool(serpapi_key=serpapi_key)
 
            presentation_agent = Agent(
                name="Web Analyst and Presenter",
                role="Expert in analysing and Searching webdata",
                goal=f"To Search Web for relevant content for information relevant to {text} and the format the data into presentation slides.",
                backstory="You are experienced web search expert who searches web and gathers data of input given and you present it to your clients.",
                tools=[search_tool],
                llm=llm,
                verbose=True
            )
 
            presentation_task_obj = Task(
                description="Get the accurate data from web searches for the input topic and give its decription and arrange it into such format that it can be given as a presentation.",
                agent=presentation_agent,
                expected_output="An factually and grammatically correct power-point presentation draft text. The text should include all the possible data and provide 10 slides with  1 topic per slide and give Decription and some noticeable facts in bullet points and  latest news on the respective topic."
            )
 
            crew = Crew(
                agents=[presentation_agent],
                tasks=[presentation_task_obj],
                verbose=True
            )
           
            result = crew.kickoff(inputs={"text": text})
            st.write(result)
           
            ppt_file = create_presentation(str(result))
            st.success("Presentation created successfully!")
           
            with open(ppt_file, "rb") as file:
                btn = st.download_button(
                    label="Download Presentation",
                    data=file,
                    file_name="presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
 
if __name__ == '__main__':
    main()
 